"""Text extraction and chunking service for document chat."""

import csv
import io

from src.core.logging import get_logger

logger = get_logger(__name__)

# Supported MIME types for extraction
SUPPORTED_TYPES = {
    "application/pdf",
    "text/plain",
    "text/markdown",
    "text/csv",
    "text/x-github-repo",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class TextExtractor:
    """Extract text from various file types and chunk for embedding."""

    def extract(self, file_path: str, content_type: str) -> str:
        """Extract text content from a file based on its content type."""
        if content_type == "application/pdf":
            return self._extract_pdf(file_path)
        elif content_type in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ):
            return self._extract_docx(file_path)
        elif content_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ):
            return self._extract_pptx(file_path)
        elif content_type == "text/csv":
            return self._extract_csv(file_path)
        elif content_type in ("text/plain", "text/markdown", "text/x-github-repo"):
            return self._extract_text(file_path)
        else:
            raise ValueError(f"Unsupported content type: {content_type}")

    def chunk_text(
        self, text: str, chunk_size: int = 1000, overlap: int = 200
    ) -> list[str]:
        """Split text into overlapping chunks for embedding."""
        if not text or not text.strip():
            return []

        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = start + chunk_size

            # Try to break at sentence boundary
            if end < text_len:
                # Look for sentence end within last 20% of chunk
                search_start = end - int(chunk_size * 0.2)
                last_period = text.rfind(". ", search_start, end)
                last_newline = text.rfind("\n", search_start, end)
                break_point = max(last_period, last_newline)
                if break_point > search_start:
                    end = break_point + 1

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - overlap
            if start >= text_len:
                break

        return chunks

    def _extract_pdf(self, file_path: str) -> str:
        import pdfplumber

        pages = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
        return "\n\n".join(pages)

    def _extract_docx(self, file_path: str) -> str:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    def _extract_pptx(self, file_path: str) -> str:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides_text.append("\n".join(texts))
        return "\n\n".join(slides_text)

    def _extract_csv(self, file_path: str) -> str:
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i == 0:
                    rows.append("Headers: " + " | ".join(row))
                else:
                    rows.append(" | ".join(row))
                if i >= 500:  # Limit rows to avoid huge text
                    rows.append(f"... (truncated, {i}+ rows)")
                    break
        return "\n".join(rows)

    def _extract_text(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
